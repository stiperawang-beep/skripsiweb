import sys
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def apply_beautiful_design(file_path, output_path):
    prs = Presentation(file_path)
    
    bg_color = RGBColor(245, 247, 250) # #F5F7FA
    accent_color = RGBColor(43, 108, 176) # #2B6CB0
    title_color = RGBColor(26, 54, 93) # #1A365D
    text_dark = RGBColor(45, 55, 72) # #2D3748
    
    for slide in prs.slides:
        # Try to set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add decorative bar
        left = Inches(0)
        top = Inches(0)
        width = Inches(0.25)
        height = prs.slide_height
        try:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = accent_color
            rect.line.fill.background()
            
            rect_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, prs.slide_width, Inches(0.1))
            rect_top.fill.solid()
            rect_top.fill.fore_color.rgb = accent_color
            rect_top.line.fill.background()
        except:
            pass

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            is_title = (shape == slide.shapes.title)
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Segoe UI'
                    if is_title:
                        run.font.color.rgb = title_color
                        run.font.bold = True
                    else:
                        run.font.color.rgb = text_dark

    prs.save(output_path)

if __name__ == '__main__':
    try:
        apply_beautiful_design('PPT ANNISA_REVISI.pptx', 'PPT ANNISA_PREMIUM.pptx')
    except Exception as e:
        print(f"Error: {e}")
