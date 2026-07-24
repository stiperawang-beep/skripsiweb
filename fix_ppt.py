import sys
from pptx import Presentation

def replace_text_in_runs(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        old_text = paragraph.text
        new_text = old_text
        for old, new in replacements.items():
            new_text = new_text.replace(old, new)
            
        if old_text != new_text:
            for run in paragraph.runs:
                run_text = run.text
                for old, new in replacements.items():
                    run_text = run_text.replace(old, new)
                run.text = run_text

            if paragraph.text != new_text:
                paragraph.text = new_text

def process_presentation(file_path, output_path):
    prs = Presentation(file_path)
    
    replacements = {
        "bebagai": "berbagai",
        "Langkat Siak.": "Langkat Siak?",
        "Langjkat": "Langkat",
        "Saik.": "Siak.",
        "mengenai ilmiah": "ilmiah",
        "tatap muka publisitas": "tatap muka, publisitas",
        "merk": "merek",
        "Tingkat kemungkinan": "tingkat kemungkinan",
        "priode": "periode",
        "ransangan": "rangsangan",
        "Kerangka Berfikir": "Kerangka Berpikir",
        "pengaruh kompensasi t erhadap kinerja karyawan": "pengaruh promosi terhadap minat beli konsumen",
        "pengaruh  kompensasi   t erhadap  kinerja karyawan": "pengaruh promosi terhadap minat beli konsumen",
        "Berdasarkanlandasan": "Berdasarkan landasan",
        "dapatdisusun": "dapat disusun",
        "rumah makan langkat": "Rumah Makan Langkat",
        "Provinsi Riau. .": "Provinsi Riau. ",
        "tahap penyusunan skripsi,": "tahap penyusunan skripsi.",
        "perhari": "per hari",
        "lemeshow": "Lemeshow",
        "konsumen, Berdasarkan": "konsumen. Berdasarkan",
        "Total_X": "Promosi (X)",
        "konsumen,  Saran": "konsumen. Saran",
        "THANK YOU FOR ATTENTION": "THANK YOU FOR YOUR ATTENTION",
        "ôBagaimana": "\"Bagaimana",
        "Siak.ö": "Siak?\""
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_runs(cell, replacements)
            else:
                replace_text_in_runs(shape, replacements)

    prs.save(output_path)

if __name__ == '__main__':
    process_presentation('PPT ANNISA.pptx', 'PPT ANNISA_REVISI.pptx')
