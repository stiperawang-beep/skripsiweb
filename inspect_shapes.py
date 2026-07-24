from pptx import Presentation

prs = Presentation('PPT ANNISA.pptx')
for i, slide in enumerate(prs.slides):
    print(f"--- SLIDE {i+1} ---")
    for shape in slide.shapes:
        shape_type = ""
        if shape.has_text_frame:
            shape_type = "TextFrame"
        elif shape.has_table:
            shape_type = "Table"
        elif shape.has_chart:
            shape_type = "Chart"
        elif hasattr(shape, 'image'):
            shape_type = "Image"
        elif shape.shape_type == 14: # Placeholder
            shape_type = "Placeholder"
        elif shape.shape_type == 13: # Picture
            shape_type = "Picture"
        elif shape.shape_type == 6: # Group
            shape_type = "Group"
        else:
            shape_type = f"Other (type: {shape.shape_type})"
        
        print(f"- Shape: {shape_type} | Name: {shape.name}")
