from pptx import Presentation
import os

os.makedirs('assets', exist_ok=True)
prs = Presentation('PPT ANNISA.pptx')

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    pic_num = 1
    for shape in slide.shapes:
        if shape.shape_type == 13: # Picture
            image = shape.image
            image_bytes = image.blob
            ext = image.ext
            filename = f"assets/slide{slide_num}_pic{pic_num}.{ext}"
            with open(filename, 'wb') as f:
                f.write(image_bytes)
            print(f"Extracted {filename}")
            pic_num += 1
